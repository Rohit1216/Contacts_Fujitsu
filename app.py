"""
LinkedIn Team Deck Generator - Streamlit UI

Upload a PowerPoint template (with one or more "team grid" slides and,
optionally, a contact-table slide) plus an Excel workbook (one sheet per
team, columns Name/Phone/Email/Designation/Location/PhotoUrl/LinkedinUrl),
map sheets to template slides, and download a fully populated deck.

Run:
    streamlit run app.py
"""

import io

import streamlit as st
from pptx import Presentation

from utils.excel_loader import load_workbook_people, dedupe_by_linkedin
from utils.ppt_generator import (
    DeckTheme,
    find_person_slots,
    paginate_team_sheet,
    paginate_contact_table,
    _table_shape,
)

st.set_page_config(page_title="LinkedIn Team Deck Generator", layout="wide")
st.title("🔗 LinkedIn Team Deck Generator")
st.caption(
    "Fill a branded PowerPoint org-chart/team-card template from an Excel "
    "workbook — photos, names (hyperlinked to LinkedIn), and titles — with "
    "automatic pagination when a team has more people than the template has slots."
)

with st.sidebar:
    st.header("1. Upload files")
    template_file = st.file_uploader("PowerPoint template (.pptx)", type=["pptx"])
    workbook_file = st.file_uploader("Team workbook (.xlsx)", type=["xlsx"])

    st.header("2. Options")
    photo_size = st.slider("Photo render size (px)", 200, 800, 400, step=50)
    max_retries = st.slider("Photo download retries", 1, 5, 3)

if not (template_file and workbook_file):
    st.info("Upload a template and a workbook in the sidebar to get started.")
    st.stop()

# Cache uploads across reruns within a session
template_bytes = template_file.getvalue()
prs_preview = Presentation(io.BytesIO(template_bytes))
people_by_sheet = load_workbook_people(workbook_file)

slide_labels = []
for i, slide in enumerate(prs_preview.slides):
    n_slots = len(find_person_slots(slide))
    if n_slots:
        slide_labels.append(f"Slide {i + 1} — team grid ({n_slots} slots)")
    elif _table_shape(slide) is not None:
        slide_labels.append(f"Slide {i + 1} — contact table")
    else:
        slide_labels.append(f"Slide {i + 1} — other")

st.header("3. Map Excel sheets → template slides")
st.caption(
    "Each 'team grid' slide is auto-paginated if a sheet has more people than "
    "it has slots. Pick one slide as the contact table (optional)."
)

sheet_to_slide = {}
cols = st.columns(2)
for idx, sheet_name in enumerate(people_by_sheet.keys()):
    with cols[idx % 2]:
        choice = st.selectbox(
            f"'{sheet_name}' ({len(people_by_sheet[sheet_name])} people) →",
            options=["(skip)"] + slide_labels,
            key=f"map_{sheet_name}",
        )
        if choice != "(skip)":
            sheet_to_slide[sheet_name] = slide_labels.index(choice)

table_choice = st.selectbox(
    "Contact table slide (combines all sheets, deduped) →",
    options=["(none)"] + slide_labels,
    index=0,
)

st.divider()

if st.button("Generate deck", type="primary"):
    prs = Presentation(io.BytesIO(template_bytes))
    original_slides = list(prs.slides)
    theme = DeckTheme(photo_size_px=photo_size, max_retries=max_retries)
    photo_cache = {}

    progress = st.progress(0.0, text="Starting...")
    status = st.empty()

    def make_cb(section_weight, base):
        def cb(done, total, message):
            frac = base + section_weight * (done / max(total, 1))
            progress.progress(min(frac, 1.0), text=message)

        return cb

    n_sections = len(sheet_to_slide) + (1 if table_choice != "(none)" else 0)
    weight = 1.0 / max(n_sections, 1)
    base = 0.0

    for sheet_name, slide_idx in sheet_to_slide.items():
        people = people_by_sheet[sheet_name]
        status.write(f"**{sheet_name}** — {len(people)} people")
        paginate_team_sheet(
            prs,
            original_slides[slide_idx],
            people,
            theme,
            photo_cache,
            title_prefix=sheet_name,
            progress_cb=make_cb(weight, base),
        )
        base += weight

    if table_choice != "(none)":
        table_idx = slide_labels.index(table_choice)
        all_people = dedupe_by_linkedin(*people_by_sheet.values())
        status.write(f"**Contact table** — {len(all_people)} unique people")
        paginate_contact_table(
            prs,
            original_slides[table_idx],
            all_people,
            theme,
            title_prefix="Contact Details of the Account Stakeholders",
        )
        progress.progress(1.0, text="Done")

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)

    st.success("Deck generated!")
    st.download_button(
        "⬇️ Download generated .pptx",
        data=out,
        file_name="Team_Deck_Generated.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
