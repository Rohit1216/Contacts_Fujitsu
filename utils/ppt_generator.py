"""
LinkedIn Team Deck Generator - core engine.

Reads an Excel workbook (one sheet per team, columns:
Name | Phone | Email | Designation | Location | PhotoUrl | LinkedinUrl)
and fills a PowerPoint template that contains:
  - one or more "team grid" slides (a big top card + a grid of small
    cards, each card = background rectangle + circular photo + hyperlinked
    name + designation), which get auto-duplicated/paginated if a sheet
    has more people than the template has slots.
  - one or more "contact table" slides (S.No | Name | Title | Phone | Email),
    which also get auto-duplicated/paginated as rows run out.

This module has no Streamlit dependency - it's plain python-pptx / PIL /
requests so it can be unit tested or reused from a CLI.
"""

import copy
import io
import re
from dataclasses import dataclass, field
from typing import Optional

import requests
from PIL import Image, ImageDraw, ImageOps
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# --------------------------------------------------------------------------
# Theme / tunables
# --------------------------------------------------------------------------


@dataclass
class DeckTheme:
    """Small config surface so callers don't have to touch internals."""

    photo_size_px: int = 400          # circular photo render resolution
    request_timeout: int = 12
    max_retries: int = 3
    placeholder_initials_bg: str = "8C1D6E"   # used when a photo can't be fetched
    rows_per_table_slide: Optional[int] = None  # None = auto-detect from template


# --------------------------------------------------------------------------
# Photo handling
# --------------------------------------------------------------------------

_session = requests.Session()
_session.headers.update(
    {"User-Agent": "Mozilla/5.0 (compatible; TeamDeckGenerator/1.0)"}
)


def _initials(name: str) -> str:
    parts = [p for p in re.split(r"\s+", name.strip()) if p]
    if not parts:
        return "?"
    if len(parts) == 1:
        return parts[0][0].upper()
    return (parts[0][0] + parts[-1][0]).upper()


def _placeholder_avatar(name: str, size: int, bg_hex: str) -> Image.Image:
    """Circular initials avatar used when a photo URL is missing/broken."""
    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    bg = tuple(int(bg_hex[i : i + 2], 16) for i in (0, 2, 4)) + (255,)
    draw.ellipse((0, 0, size, size), fill=bg)
    text = _initials(name)
    from PIL import ImageFont

    try:
        font = ImageFont.truetype(
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", size // 2
        )
    except Exception:
        font = ImageFont.load_default()
    bbox = draw.textbbox((0, 0), text, font=font)
    w, h = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text(
        ((size - w) / 2 - bbox[0], (size - h) / 2 - bbox[1]),
        text,
        font=font,
        fill=(255, 255, 255, 255),
    )
    return img


def fetch_circular_photo(
    url: Optional[str], name: str, theme: DeckTheme, cache: Optional[dict] = None
) -> io.BytesIO:
    """
    Download `url`, crop to a square, mask to a circle, return PNG bytes.
    Falls back to an initials avatar if the URL is empty or the download fails.
    Results are memoized in `cache` (keyed by url) so repeated people/pages
    only hit the network once.
    """
    size = theme.photo_size_px
    if cache is not None and url and url in cache:
        return io.BytesIO(cache[url])

    img = None
    if url:
        for attempt in range(theme.max_retries):
            try:
                resp = _session.get(url, timeout=theme.request_timeout)
                resp.raise_for_status()
                img = Image.open(io.BytesIO(resp.content)).convert("RGBA")
                break
            except Exception:
                img = None
                continue

    if img is None:
        img = _placeholder_avatar(name, size, theme.placeholder_initials_bg)
    else:
        # center-crop to square, then resize
        w, h = img.size
        side = min(w, h)
        left, top = (w - side) // 2, (h - side) // 2
        img = img.crop((left, top, left + side, top + side)).resize(
            (size, size), Image.LANCZOS
        )

    # circular alpha mask
    mask = Image.new("L", (size, size), 0)
    ImageDraw.Draw(mask).ellipse((0, 0, size, size), fill=255)
    out = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    out.paste(img, (0, 0), mask)

    buf = io.BytesIO()
    out.save(buf, format="PNG")
    data = buf.getvalue()
    if cache is not None and url:
        cache[url] = data
    return io.BytesIO(data)


# --------------------------------------------------------------------------
# Slide/shape introspection - find the repeating "person card" slots
# --------------------------------------------------------------------------


@dataclass
class PersonSlot:
    rect: object
    picture: object
    name_box: object
    desig_box: object


def find_person_slots(slide) -> list:
    """
    Auto-detect the repeating (photo + name + designation [+ background
    rectangle]) groups on a template slide, in reading order (top-to-bottom,
    left-to-right). Works regardless of shape id/z-order quirks because it
    matches purely on geometry.
    """
    pictures = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
    textboxes = [s for s in slide.shapes if s.shape_type == 17]  # TEXT_BOX only
    rects = [s for s in slide.shapes if s.shape_type == 1]  # AUTO_SHAPE (cards)

    slots = []
    for pic in pictures:
        # the two textboxes that sit to the right of this picture, vertically
        # within its band
        candidates = [
            tb
            for tb in textboxes
            if 0 <= (tb.left - (pic.left + pic.width)) <= pic.width * 3
            and tb.top >= pic.top - pic.height * 0.5
            and tb.top <= pic.top + pic.height * 1.5
        ]
        candidates.sort(key=lambda tb: tb.top)
        if len(candidates) < 2:
            continue
        name_box, desig_box = candidates[0], candidates[1]

        # background rectangle whose bounds contain the picture
        rect = None
        for r in rects:
            if (
                r.left <= pic.left
                and r.top <= pic.top
                and r.left + r.width >= pic.left + pic.width
                and r.top + r.height >= pic.top + pic.height
            ):
                rect = r
                break

        slots.append(PersonSlot(rect, pic, name_box, desig_box))

    def row_key(slot):
        return (round(slot.picture.top / 50000), slot.picture.left)

    slots.sort(key=row_key)
    return slots


# --------------------------------------------------------------------------
# Low-level text / hyperlink / picture replacement helpers
# --------------------------------------------------------------------------


def _set_run_text(text_frame, text: str):
    """Set text on the first run of the first paragraph, preserving its
    formatting (font, color, bold, etc). Extra paragraphs/runs are removed."""
    paragraphs = text_frame.paragraphs
    first_p = paragraphs[0]
    if not first_p.runs:
        first_p.add_run()
    first_run = first_p.runs[0]
    first_run.text = text
    # drop any additional runs in the first paragraph
    for extra in first_p.runs[1:]:
        extra._r.getparent().remove(extra._r)
    # drop any additional paragraphs
    for p in paragraphs[1:]:
        p._p.getparent().remove(p._p)
    return first_run


def _autofit_shrink(run, box_width_emu: int, min_pt: float = 7.0):
    """Long real-world names/titles routinely overflow the narrow template
    boxes that were sized for placeholder 'xxxxx' text. Rather than let
    PowerPoint clip the text, shrink the run's font size just enough to fit
    a single line (heuristic: no font metrics available at this layer)."""
    text = run.text
    if not text or not box_width_emu:
        return
    size = run.font.size
    if size is None:
        return
    pt = size.pt
    bold = bool(run.font.bold)
    per_char_em = 0.62 if bold else 0.52
    avg_char_w_emu = pt * per_char_em * 12700
    needed = avg_char_w_emu * len(text)
    # small safety margin - these boxes have zero internal padding
    if needed > box_width_emu * 0.94:
        scale = (box_width_emu * 0.94) / needed
        new_pt = max(min_pt, pt * scale)
        run.font.size = Pt(round(new_pt, 1))


def set_name_with_hyperlink(name_box, name: str, linkedin_url: Optional[str]):
    run = _set_run_text(name_box.text_frame, name)
    rPr = run._r.get_or_add_rPr()
    existing = rPr.find(qn("a:hlinkClick"))
    if linkedin_url:
        run.hyperlink.address = linkedin_url
    elif existing is not None:
        rPr.remove(existing)
    _autofit_shrink(run, name_box.width)


def set_designation(desig_box, text: str):
    run = _set_run_text(desig_box.text_frame, text or "")
    _autofit_shrink(run, desig_box.width)


def replace_picture(slide, pic_shape, image_stream):
    """Swap the image inside an existing picture placeholder while keeping
    its exact position/size."""
    left, top, width, height = (
        pic_shape.left,
        pic_shape.top,
        pic_shape.width,
        pic_shape.height,
    )
    old_el = pic_shape._element
    parent = old_el.getparent()
    idx = list(parent).index(old_el)
    parent.remove(old_el)
    new_pic = slide.shapes.add_picture(image_stream, left, top, width, height)
    new_el = new_pic._element
    parent.remove(new_el)
    parent.insert(idx, new_el)
    return new_pic


def remove_slot(slide, slot: PersonSlot):
    """Delete an entire unused person card (background + photo + text)."""
    for shape in (slot.rect, slot.picture, slot.name_box, slot.desig_box):
        if shape is None:
            continue
        el = shape._element
        p = el.getparent()
        if p is not None:
            p.remove(el)


# --------------------------------------------------------------------------
# Slide duplication (python-pptx has no native API for this)
# --------------------------------------------------------------------------


def duplicate_slide(prs: Presentation, source_slide, insert_after_slide=None):
    """
    Clone `source_slide` (including images + hyperlink relationships),
    append it to the deck, then move it to sit directly after
    `insert_after_slide` (defaults to right after the source).
    Returns the new slide.
    """
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # remove any placeholder shapes python-pptx auto-added from the layout
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    # deep-copy every shape from the source
    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(new_el)

    # copy over every relationship the source slide uses (images, hyperlinks,
    # etc). python-pptx assigns fresh rIds, so build an old->new map and
    # rewrite the r:id/r:embed references inside the copied shape XML.
    rid_map = {}
    for rel_id, rel in source_slide.part.rels.items():
        if rel.reltype.endswith("/slideLayout") or rel.reltype.endswith(
            "/notesSlide"
        ):
            continue
        if rel.is_external:
            new_rid = new_slide.part.rels._add_relationship(
                rel.reltype, rel.target_ref, is_external=True
            )
        else:
            new_rid = new_slide.part.rels._add_relationship(
                rel.reltype, rel.target_part
            )
        rid_map[rel_id] = new_rid

    R_ID_ATTRS = {qn("r:id"), qn("r:embed"), qn("r:link")}
    for shape in new_slide.shapes:
        for el in shape._element.iter():
            for attr in R_ID_ATTRS:
                old_val = el.attrib.get(attr)
                if old_val and old_val in rid_map:
                    el.attrib[attr] = rid_map[old_val]

    # reposition in the slide order
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    new_pos_el = slides[-1]  # the one we just appended, currently last
    xml_slides.remove(new_pos_el)

    anchor = insert_after_slide or source_slide
    anchor_idx = list(prs.slides).index(anchor)
    xml_slides.insert(anchor_idx + 1, new_pos_el)

    return new_slide


# --------------------------------------------------------------------------
# High level: fill one "team grid" page
# --------------------------------------------------------------------------


def fill_team_page(
    slide, people: list, theme: DeckTheme, photo_cache: dict, progress_cb=None
):
    """
    people: list of dicts with keys name, designation, photo_url, linkedin_url
            length must be <= number of slots on the slide.
    Any unused slots (when this is the last, partial page) are deleted.
    progress_cb(done, total, message): optional callback for UI progress bars.
    """
    slots = find_person_slots(slide)
    for i, slot in enumerate(slots):
        if i >= len(people):
            remove_slot(slide, slot)
            continue
        person = people[i]
        if progress_cb:
            progress_cb(i, len(people), f"Adding {person['name']}...")
        set_name_with_hyperlink(slot.name_box, person["name"], person.get("linkedin_url"))
        set_designation(slot.desig_box, person.get("designation", ""))
        photo = fetch_circular_photo(
            person.get("photo_url"), person["name"], theme, photo_cache
        )
        replace_picture(slide, slot.picture, photo)
    if progress_cb:
        progress_cb(len(people), len(people), "Page complete")


def paginate_team_sheet(
    prs: Presentation,
    template_slide,
    people: list,
    theme: DeckTheme,
    photo_cache: dict,
    title_prefix: Optional[str] = None,
    progress_cb=None,
):
    """Fill `template_slide` with the first chunk of `people`, then duplicate
    it as many times as needed for the rest. Updates the "(x/y)" style title
    if a title textbox is found and title_prefix is given."""
    slot_count = len(find_person_slots(template_slide))
    if slot_count == 0:
        raise ValueError("No person slots detected on template slide.")

    chunks = [
        people[i : i + slot_count] for i in range(0, max(len(people), 1), slot_count)
    ] or [[]]
    total_pages = len(chunks)

    pages = [template_slide]
    prev = template_slide
    for _ in chunks[1:]:
        new_slide = duplicate_slide(prs, template_slide, insert_after_slide=prev)
        pages.append(new_slide)
        prev = new_slide

    for page_idx, (slide, chunk) in enumerate(zip(pages, chunks), start=1):
        fill_team_page(slide, chunk, theme, photo_cache, progress_cb=progress_cb)
        if title_prefix:
            _update_page_title(slide, title_prefix, page_idx, total_pages)

    return pages


def _update_page_title(slide, prefix: str, page: int, total: int):
    """Update just the '(x/y)' page-count fragment of the title, in place,
    so the rest of the title's text/formatting is untouched. Falls back to
    replacing the whole title only if no such fragment is found."""
    pattern = re.compile(r"\(\d+\s*/\s*\d+\)")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if pattern.fullmatch(run.text.strip()):
                    run.text = f"({page}/{total})"
                    return
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.lower().startswith("title"):
            _set_run_text(shape.text_frame, f"{prefix} ({page}/{total})")
            return


# --------------------------------------------------------------------------
# High level: contact table pages
# --------------------------------------------------------------------------


def _table_shape(slide):
    for shape in slide.shapes:
        if shape.has_table:
            return shape
    return None


def paginate_contact_table(
    prs: Presentation,
    template_slide,
    people: list,
    theme: DeckTheme,
    title_prefix: Optional[str] = None,
):
    table_shape = _table_shape(template_slide)
    if table_shape is None:
        raise ValueError("No table found on contact-table template slide.")
    table = table_shape.table
    header_row_count = 1
    data_rows_in_template = len(table.rows) - header_row_count
    rows_per_page = theme.rows_per_table_slide or data_rows_in_template

    chunks = [
        people[i : i + rows_per_page]
        for i in range(0, max(len(people), 1), rows_per_page)
    ] or [[]]

    pages = [template_slide]
    prev = template_slide
    for _ in chunks[1:]:
        new_slide = duplicate_slide(prs, template_slide, insert_after_slide=prev)
        pages.append(new_slide)
        prev = new_slide

    total_pages = len(chunks)
    running_index = 1
    for page_idx, (slide, chunk) in enumerate(zip(pages, chunks), start=1):
        t_shape = _table_shape(slide)
        _fill_table_page(t_shape.table, chunk, rows_per_page, running_index)
        running_index += len(chunk)
        if title_prefix:
            _update_page_title(slide, title_prefix, page_idx, total_pages)

    return pages


def _fill_table_page(table, people_chunk, rows_needed, start_index):
    header_row = table.rows[0]
    existing_data_rows = len(table.rows) - 1

    # grow the table if this page needs more rows than the template has
    tbl = table._tbl
    while len(table.rows) - 1 < rows_needed:
        last_tr = tbl.tr_lst[-1]
        new_tr = copy.deepcopy(last_tr)
        tbl.append(new_tr)

    # shrink (remove trailing rows) if fewer are needed on the final page
    while len(table.rows) - 1 > max(rows_needed, len(people_chunk)):
        tbl.remove(tbl.tr_lst[-1])

    for i in range(len(table.rows) - 1):
        row = table.rows[i + 1]
        if i < len(people_chunk):
            p = people_chunk[i]
            values = [
                str(start_index + i),
                p["name"],
                p.get("designation", ""),
                p.get("phone", "") or "-",
                p.get("email", "") or "-",
            ]
            for cell, val in zip(row.cells, values):
                _set_run_text(cell.text_frame, val)
        else:
            # blank leftover row on a partial final page
            for cell in row.cells:
                _set_run_text(cell.text_frame, "")
